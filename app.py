from flask import Flask, request, jsonify, send_file, redirect, url_for, render_template, session, flash
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import requests
import os
import json
import io
import tempfile
import uuid
import logging
from template_manager import TemplateManager
from werkzeug.security import generate_password_hash, check_password_hash
import sqlite3
from io import BytesIO
from flask_cors import CORS
from PIL import Image
from functools import wraps

app = Flask(__name__)
CORS(app)
template_manager = TemplateManager()

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

app.secret_key = 'your_secret_key'  # Change this to a secure random string in production

def get_db():
    conn = sqlite3.connect('users.db')
    conn.row_factory = sqlite3.Row
    return conn

# Initialize the database if it doesn't exist
def init_db():
    conn = get_db()
    cursor = conn.cursor()
    # Create users table
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        email TEXT UNIQUE NOT NULL,
        password TEXT NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    ''')
    
    # Create presentations table to store user's presentation history
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS presentations (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        title TEXT NOT NULL,
        filename TEXT NOT NULL,
        template TEXT NOT NULL,
        slide_count INTEGER NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users (id)
    )
    ''')
    
    conn.commit()
    conn.close()
    logger.info("Database initialized with required tables.")

# Login decorator
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

@app.route('/register', methods=['GET', 'POST'])
def register():
    if 'user_id' in session:
        return redirect(url_for('dashboard'))
        
    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']
        confirm_password = request.form['confirm_password']
        
        # Validate input
        if not username or not email or not password:
            return render_template('register.html', error='All fields are required')
            
        if password != confirm_password:
            return render_template('register.html', error='Passwords do not match')
            
        if len(password) < 6:
            return render_template('register.html', error='Password must be at least 6 characters')
        
        hashed_password = generate_password_hash(password)
        
        try:
            conn = get_db()
            c = conn.cursor()
            
            # Check if username or email already exists
            c.execute('SELECT * FROM users WHERE username = ? OR email = ?', (username, email))
            existing_user = c.fetchone()
            
            if existing_user:
                conn.close()
                return render_template('register.html', error='Username or email already exists')
                
            # Insert new user
            c.execute('INSERT INTO users (username, email, password) VALUES (?, ?, ?)',
                    (username, email, hashed_password))
            conn.commit()
            
            # Get the user id for session
            c.execute('SELECT id FROM users WHERE email = ?', (email,))
            user_id = c.fetchone()[0]
            conn.close()
            
            # Set session data
            session['user_id'] = user_id
            session['username'] = username
            
            return redirect(url_for('dashboard'))
            
        except sqlite3.Error as e:
            logger.error(f"Database error during registration: {str(e)}")
            return render_template('register.html', error='An error occurred. Please try again.')
            
    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if 'user_id' in session:
        return redirect(url_for('dashboard'))
        
    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        
        try:
            conn = get_db()
            c = conn.cursor()
            c.execute('SELECT * FROM users WHERE email = ?', (email,))
            user = c.fetchone()
            conn.close()
            
            if user and check_password_hash(user['password'], password):
                session['user_id'] = user['id']
                session['username'] = user['username']
                return redirect(url_for('dashboard'))
            else:
                return render_template('login.html', error='Invalid email or password')
                
        except sqlite3.Error as e:
            logger.error(f"Database error during login: {str(e)}")
            return render_template('login.html', error='An error occurred. Please try again.')
            
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.pop('user_id', None)
    session.pop('username', None)
    return redirect(url_for('login'))

from datetime import datetime

@app.route('/dashboard')
@login_required
def dashboard():
    # Get user's presentation history
    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT * FROM presentations WHERE user_id = ? ORDER BY created_at DESC', (session['user_id'],))
    presentations = c.fetchall()
    conn.close()

    # Convert 'created_at' from string to datetime object
    presentation_list = []
    for pres in presentations:
        # Convert sqlite3.Row to dictionary
        pres_dict = dict(pres)  # Convert to a dictionary for mutability

        # If 'created_at' is a string, convert it to a datetime object
        if isinstance(pres_dict['created_at'], str):
            pres_dict['created_at'] = datetime.strptime(pres_dict['created_at'], "%Y-%m-%d %H:%M:%S")
        
        # Add the modified presentation to the new list
        presentation_list.append(pres_dict)

    # Get all available templates
    templates = template_manager.get_all_templates()
    template_list = []
    for key, template in templates.items():
        template_list.append({
            "id": key,
            "name": template.get('name', key),
            "description": template.get('description', ''),
            "preview_image": template.get('preview_image', '')
        })
    
    return render_template('dashboard.html', 
                          username=session['username'], 
                          presentations=presentation_list,
                          templates=template_list)

OLLAMA_ENDPOINT = "http://localhost:11434/api/generate"

def generate_text_content(topic, num_slides, custom_content=None):
    try:
        if custom_content:
            # If custom content is provided, ask Ollama to format it
            prompt = f"""Convert the following presentation content into a well-structured JSON format.
            The content provided by the user is about: '{topic}'
            
            USER CONTENT:
            {custom_content}
            
            Format EXACTLY as this JSON structure:
            {{
                "title": "Overall Presentation Title",
                "slides": [
                    {{
                        "title": "Slide 1 Title",
                        "points": [
                            "Point 1: Detailed explanation or context",
                            "Point 2: Detailed explanation or context",
                            "Point 3: Detailed explanation or context",
                            "Point 4: Additional context or related points",
                            "Point 5: Further insights or examples"
                        ]
                    }},
                    ...
                ]
            }}
            
            Requirements:
            - Use clear, professional language
            - Extract slide titles and bullet points from the user content
            - Organize the content logically
            - If the user hasn't provided enough structure, create appropriate slide titles and organize the content
            - Add bullet points where not explicitly provided by user
            - Avoid any markdown, code blocks, or extra formatting
            """
        else:
            # Original prompt for generating from a topic
            prompt = f"""Generate a detailed JSON for a presentation about '{topic}' with {num_slides} slides.
            Each slide should have the following:
            - A detailed title
            - At least 5 concise and informative bullet points per slide (if applicable)
            - Provide some additional explanations or insights for each bullet point
            - Ensure the content is rich, professional, and informative
            Format EXACTLY as this JSON structure:
            {{
                "title": "Overall Presentation Title",
                "slides": [
                    {{
                        "title": "Slide 1 Title",
                        "points": [
                            "Point 1: Detailed explanation or context",
                            "Point 2: Detailed explanation or context",
                            "Point 3: Detailed explanation or context",
                            "Point 4: Additional context or related points",
                            "Point 5: Further insights or examples"
                        ]
                    }},
                    ...
                ]
            }}
            Requirements:
            - Use clear, professional language
            - Ensure each slide has a meaningful title
            - Create at least 5 detailed, informative bullet points per slide
            - Provide explanations, context, or examples where relevant
            - Avoid any markdown, code blocks, or extra formatting
            """
        
        payload = {
            "model": "gemma3:1b-it-qat",
            "prompt": prompt,
            "stream": False,
            "format": "json"
        }
        response = requests.post(OLLAMA_ENDPOINT, json=payload)
        if response.status_code != 200:
            logger.error(f"Ollama API error: {response.status_code} - {response.text}")
            raise Exception(f"Ollama API error: {response.status_code}")
        content = response.json()["response"]
        if "```json" in content:
            content = content.split("```json")[1].split("```")[0].strip()
        elif "```" in content:
            content = content.split("```")[1].split("```")[0].strip()
        presentation_data = json.loads(content)
        if not isinstance(presentation_data, dict) or 'title' not in presentation_data or 'slides' not in presentation_data:
            raise ValueError("Invalid JSON structure")
        for slide in presentation_data.get('slides', []):
            if 'title' not in slide or 'points' not in slide:
                raise ValueError("Invalid slide structure")
        return presentation_data
    except Exception as e:
        logger.error(f"Text generation error: {str(e)}")
        # Create a fallback presentation structure
        title = topic or "Presentation"
        return {
            "title": title,
            "slides": [
                {
                    "title": f"Introduction to {title}",
                    "points": [
                        "Overview of the topic with more context and background",
                        "Key points to discuss with additional details",
                        "Importance and relevance with examples or data"
                    ]
                },
                {
                    "title": "Main Concepts",
                    "points": [
                        "First main concept with detailed examples",
                        "Second main concept with further elaboration",
                        "Third main concept with supporting data or case studies"
                    ]
                },
                {
                    "title": "Conclusion",
                    "points": [
                        "Summary of key takeaways with insights",
                        "Future implications with potential applications",
                        "Call to action with a proposed next step or idea"
                    ]
                }
            ]
        }

def generate_image_prompt(prompt):
    return f"Professional presentation image related to: {prompt}"

def create_presentation(content_data, image_prompts=None, template="default"):
    try:
        template_config = template_manager.get_template(template) or template_manager.get_template('default')
        styles = template_config.get('styles', {})
        title_slide_styles = styles.get('title_slide', {})
        content_slide_styles = styles.get('content_slide', {})
        image_slide_styles = styles.get('image_slide', {})
        
        # Log preview image
        preview_image = template_config.get('preview_image', '')
        preview_image_path = os.path.join('static', preview_image)
        if preview_image and os.path.exists(preview_image_path):
            logger.info(f"Preview image found for template {template}: {preview_image_path}")
        else:
            logger.warning(f"Preview image not found for template {template}: {preview_image_path}")
        
        preview_data = {
            "title": content_data.get("title", "Presentation"),
            "template": template,
            "styles": {
                "title_slide": {
                    "background": title_slide_styles.get('background', {'type': 'solid', 'color': {'r': 240, 'g': 240, 'b': 240}}),
                    "background_image": title_slide_styles.get('background_image', ''),
                    "title_font": title_slide_styles.get('title_font', {'name': 'Calibri', 'size': 44, 'color': {'r': 0, 'g': 0, 'b': 0}, 'bold': True, 'alignment': 'center'}),
                    "image_position": title_slide_styles.get('image_position', {'left': 2.5, 'top': 4.0, 'width': 5.0, 'height': 2.5})
                },
                "content_slide": {
                    "background": content_slide_styles.get('background', {'type': 'solid', 'color': {'r': 255, 'g': 255, 'b': 255}}),
                    "background_image": content_slide_styles.get('background_image', ''),
                    "title_font": content_slide_styles.get('title_font', {'name': 'Calibri', 'size': 32, 'color': {'r': 0, 'g': 0, 'b': 0}, 'bold': True, 'alignment': 'left'}),
                    "body_font": content_slide_styles.get('body_font', {'name': 'Calibri', 'size': 18, 'color': {'r': 50, 'g': 50, 'b': 50}, 'alignment': 'left'}),
                    "image_position": content_slide_styles.get('image_position', {'left': 6.0, 'top': 1.5, 'width': 3.5, 'height': 4.5})
                },
                "image_slide": {
                    "fill_color": image_slide_styles.get('fill_color', {'r': 245, 'g': 245, 'b': 245}),
                    "border_color": image_slide_styles.get('border_color', {'r': 200, 'g': 200, 'b': 200}),
                    "border_width": image_slide_styles.get('border_width', 1.5),
                    "border_style": image_slide_styles.get('border_style', 'dashed')
                }
            },
            "slides": []
        }
        
        prs = Presentation()
        SLIDE_WIDTH = Inches(10)  # Standard 4:3 slide width
        SLIDE_HEIGHT = Inches(7.5)  # Standard 4:3 slide height
        SUPPORTED_FORMATS = {'BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF'}
        
        def validate_image_format(image_path):
            try:
                with Image.open(image_path) as img:
                    format = img.format.upper()
                    if format not in SUPPORTED_FORMATS:
                        logger.warning(f"Unsupported image format at {image_path}: got {format}, expected one of {SUPPORTED_FORMATS}")
                        return False
                    logger.debug(f"Validated image format at {image_path}: {format}")
                    return True
            except Exception as e:
                logger.error(f"Failed to validate image format at {image_path}: {str(e)}")
                return False
        
        def adjust_font_size(title_text, base_size):
            # Reduce font size for long titles (e.g., > 40 characters)
            if len(title_text) > 40:
                new_size = max(base_size - 8, 20)  # Reduce by up to 8pt, minimum 20pt
                logger.debug(f"Reducing font size for title '{title_text[:20]}...': {base_size}pt to {new_size}pt")
                return new_size
            return base_size
        
        # Title Slide
        blank_slide_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(blank_slide_layout)
        
        background_settings = title_slide_styles.get('background', {})
        bg_image = title_slide_styles.get('background_image', '')
        background = title_slide.background
        fill = background.fill
        if bg_image:
            bg_image_path = os.path.abspath(os.path.join('static', bg_image))
            logger.debug(f"Checking background image for title slide: {bg_image_path}")
            if os.path.exists(bg_image_path) and validate_image_format(bg_image_path):
                try:
                    logger.info(f"Applying background image for title slide: {bg_image_path}")
                    picture = title_slide.shapes.add_picture(
                        bg_image_path,
                        left=0,
                        top=0,
                        width=SLIDE_WIDTH,
                        height=SLIDE_HEIGHT
                    )
                    logger.debug(f"Picture added to title slide: width={picture.width.inches:.2f}in, height={picture.height.inches:.2f}in, left={picture.left.inches:.2f}in, top={picture.top.inches:.2f}in")
                    title_slide.shapes._spTree.remove(picture._element)
                    title_slide.shapes._spTree.insert(2, picture._element)
                except Exception as e:
                    logger.error(f"Failed to apply background image for title slide: {bg_image_path}, error: {str(e)}")
                    fill.solid()
                    bg_color = background_settings.get('color', {'r': 240, 'g': 240, 'b': 240}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 240, 'g': 240, 'b': 240})
                    fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                    logger.info(f"Fallback to solid color for title slide: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
            else:
                logger.error(f"Background image not found or invalid format for title slide: {bg_image_path}")
                fill.solid()
                bg_color = background_settings.get('color', {'r': 240, 'g': 240, 'b': 240}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 240, 'g': 240, 'b': 240})
                fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                logger.info(f"Fallback to solid color for title slide: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
        else:
            logger.debug(f"No background image specified for title slide, using {background_settings.get('type', 'solid')} background")
            fill.solid()
            bg_color = background_settings.get('color', {'r': 240, 'g': 240, 'b': 240}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 240, 'g': 240, 'b': 240})
            fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
            logger.info(f"Applied solid color for title slide: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
        
        # Title slide title textbox
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(2.0)  # Increased height to accommodate wrapped text
        title_box = title_slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True  # Enable word wrapping
        title_text = content_data.get("title", "Presentation")
        title_frame.text = title_text
        logger.debug(f"Title slide heading: '{title_text}', length: {len(title_text)}")
        title_para = title_frame.paragraphs[0]
        title_font_settings = title_slide_styles.get('title_font', {})
        title_para.font.name = title_font_settings.get('name', 'Calibri')
        base_font_size = title_font_settings.get('size', 44)
        title_para.font.size = Pt(adjust_font_size(title_text, base_font_size))
        title_color = title_font_settings.get('color', {'r': 0, 'g': 0, 'b': 0})
        title_para.font.color.rgb = RGBColor(title_color['r'], title_color['g'], title_color['b'])
        title_para.font.bold = title_font_settings.get('bold', True)
        title_para.alignment = PP_ALIGN.CENTER
        
        title_image_style = {}
        if image_prompts and "title" in image_prompts:
            image_position = title_slide_styles.get('image_position', {'left': 2.5, 'top': 4.0, 'width': 5.0, 'height': 2.5})
            img_left = Inches(image_position.get('left', 2.5))
            img_top = Inches(image_position.get('top', 4.0))
            img_width = Inches(image_position.get('width', 5.0))
            img_height = Inches(image_position.get('height', 2.5))
            img_placeholder = title_slide.shapes.add_shape(1, img_left, img_top, img_width, img_height)
            img_placeholder.fill.solid()
            fill_color = image_slide_styles.get('fill_color', {'r': 245, 'g': 245, 'b': 245})
            img_placeholder.fill.fore_color.rgb = RGBColor(fill_color['r'], fill_color['g'], fill_color['b'])
            border_color = image_slide_styles.get('border_color', {'r': 200, 'g': 200, 'b': 200})
            img_placeholder.line.color.rgb = RGBColor(border_color['r'], border_color['g'], border_color['b'])
            img_placeholder.line.width = Pt(image_slide_styles.get('border_width', 1.5))
            img_placeholder.line.dash_style = 2 if image_slide_styles.get('border_style', 'dashed') == 'dashed' else 1
            text_frame = img_placeholder.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            icon_p = text_frame.add_paragraph()
            icon_p.text = "🖼️"
            icon_p.alignment = PP_ALIGN.CENTER
            icon_p.font.size = Pt(48)
            icon_p.space_after = Pt(10)
            prompt_p = text_frame.add_paragraph()
            prompt_p.text = image_prompts['title']
            prompt_p.alignment = PP_ALIGN.CENTER
            prompt_p.font.italic = True
            prompt_p.font.size = Pt(14)
            prompt_p.font.color.rgb = RGBColor(100, 100, 100)
            title_image_style = {
                "left": image_position.get('left', 2.5),
                "top": image_position.get('top', 4.0),
                "width": image_position.get('width', 5.0),
                "height": image_position.get('height', 2.5),
                "fill_color": fill_color,
                "border_color": border_color,
                "border_width": image_slide_styles.get('border_width', 1.5),
                "border_style": image_slide_styles.get('border_style', 'dashed')
            }
        
        preview_data["slides"].append({
            "type": "title",
            "title": content_data.get("title", "Presentation"),
            "has_image": "title" in image_prompts if image_prompts else False,
            "image_prompt": image_prompts.get("title") if image_prompts else None,
            "image_style": title_image_style
        })
        
        # Content Slides
        for i, slide_data in enumerate(content_data.get("slides", [])):
            slide_index = str(i)
            content_slide = prs.slides.add_slide(blank_slide_layout)
            background_settings = content_slide_styles.get('background', {})
            bg_image = content_slide_styles.get('background_image', '')
            background = content_slide.background
            fill = background.fill
            if bg_image:
                bg_image_path = os.path.abspath(os.path.join('static', bg_image))
                logger.debug(f"Checking background image for content slide {i+1}: {bg_image_path}")
                if os.path.exists(bg_image_path) and validate_image_format(bg_image_path):
                    try:
                        logger.info(f"Applying background image for content slide {i+1}: {bg_image_path}")
                        picture = content_slide.shapes.add_picture(
                            bg_image_path,
                            left=0,
                            top=0,
                            width=SLIDE_WIDTH,
                            height=SLIDE_HEIGHT
                        )
                        logger.debug(f"Picture added to content slide {i+1}: width={picture.width.inches:.2f}in, height={picture.height.inches:.2f}in, left={picture.left.inches:.2f}in, top={picture.top.inches:.2f}in")
                        content_slide.shapes._spTree.remove(picture._element)
                        content_slide.shapes._spTree.insert(2, picture._element)
                    except Exception as e:
                        logger.error(f"Failed to apply background image for content slide {i+1}: {bg_image_path}, error: {str(e)}")
                        fill.solid()
                        bg_color = background_settings.get('color', {'r': 255, 'g': 255, 'b': 255}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 255, 'g': 255, 'b': 255})
                        fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                        logger.info(f"Fallback to solid color for content slide {i+1}: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
                else:
                    logger.error(f"Background image not found or invalid format for content slide {i+1}: {bg_image_path}")
                    fill.solid()
                    bg_color = background_settings.get('color', {'r': 255, 'g': 255, 'b': 255}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 255, 'g': 255, 'b': 255})
                    fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                    logger.info(f"Fallback to solid color for content slide {i+1}: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
            else:
                logger.debug(f"No background image specified for content slide {i+1}, using {background_settings.get('type', 'solid')} background")
                fill.solid()
                bg_color = background_settings.get('color', {'r': 255, 'g': 255, 'b': 255}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 255, 'g': 255, 'b': 255})
                fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                logger.info(f"Applied solid color for content slide {i+1}: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
            
            # Content slide title textbox
            title_left = Inches(0.5)
            title_top = Inches(0.5)
            title_width = Inches(9.0)
            title_height = Inches(1.2)  # Increased height for wrapped text
            title_box = content_slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True  # Enable word wrapping
            title_text = slide_data.get("title", f"Slide {i+1}")
            title_frame.text = title_text
            logger.debug(f"Content slide {i+1} heading: '{title_text}', length: {len(title_text)}")
            title_para = title_frame.paragraphs[0]
            title_font = content_slide_styles.get('title_font', {})
            title_para.font.name = title_font.get('name', 'Calibri')
            base_font_size = title_font.get('size', 32)
            title_para.font.size = Pt(adjust_font_size(title_text, base_font_size))
            title_color = title_font.get('color', {'r': 0, 'g': 0, 'b': 0})
            title_para.font.color.rgb = RGBColor(title_color['r'], title_color['g'], title_color['b'])
            title_para.font.bold = title_font.get('bold', True)
            title_para.alignment = {
                'center': PP_ALIGN.CENTER,
                'left': PP_ALIGN.LEFT,
                'right': PP_ALIGN.RIGHT
            }.get(title_font.get('alignment', 'left'), PP_ALIGN.LEFT)
            
            points_styling = []
            if slide_data.get("points", []):
                content_left = Inches(0.5)
                content_top = Inches(2.0)  # Adjusted to account for taller title
                content_width = Inches(5.0)
                # Increase height to accommodate more content
                content_height = Inches(5.0)  # Increased from 4.0 to 5.0 for more space
                content_box = content_slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Auto-fit text to shape
                body_font = content_slide_styles.get('body_font', {})
                for point in slide_data.get("points", []):
                    if text_frame.paragraphs and text_frame.paragraphs[0].text == "":
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = "• " + point
                    p.font.name = body_font.get('name', 'Calibri')
                    p.font.size = Pt(body_font.get('size', 18))
                    body_color = body_font.get('color', {'r': 50, 'g': 50, 'b': 50})
                    p.font.color.rgb = RGBColor(body_color['r'], body_color['g'], body_color['b'])
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    p.alignment = {
                        'center': PP_ALIGN.CENTER,
                        'left': PP_ALIGN.LEFT,
                        'right': PP_ALIGN.RIGHT
                    }.get(body_font.get('alignment', 'left'), PP_ALIGN.LEFT)
                    points_styling.append({
                        "text": point,
                        "level": 0,
                        "font_name": body_font.get('name', 'Calibri'),
                        "font_size": body_font.get('size', 18),
                        "color": body_color,
                        "alignment": body_font.get('alignment', 'left'),
                        "space_before": 6,
                        "space_after": 6
                    })
            
            content_image_style = {}
            if image_prompts and slide_index in image_prompts:
                image_position = content_slide_styles.get('image_position', {'left': 6.0, 'top': 2.0, 'width': 3.5, 'height': 4.0})
                img_left = Inches(image_position.get('left', 6.0))
                img_top = Inches(image_position.get('top', 2.0))
                img_width = Inches(image_position.get('width', 3.5))
                img_height = Inches(image_position.get('height', 4.0))
                img_placeholder = content_slide.shapes.add_shape(1, img_left, img_top, img_width, img_height)
                img_placeholder.fill.solid()
                fill_color = image_slide_styles.get('fill_color', {'r': 245, 'g': 245, 'b': 245})
                img_placeholder.fill.fore_color.rgb = RGBColor(fill_color['r'], fill_color['g'], fill_color['b'])
                border_color = image_slide_styles.get('border_color', {'r': 200, 'g': 200, 'b': 200})
                img_placeholder.line.color.rgb = RGBColor(border_color['r'], border_color['g'], border_color['b'])
                img_placeholder.line.width = Pt(image_slide_styles.get('border_width', 1.5))
                img_placeholder.line.dash_style = 2 if image_slide_styles.get('border_style', 'dashed') == 'dashed' else 1
                text_frame = img_placeholder.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                icon_p = text_frame.add_paragraph()
                icon_p.text = "🖼️"
                icon_p.alignment = PP_ALIGN.CENTER
                icon_p.font.size = Pt(48)
                icon_p.space_after = Pt(10)
                prompt_p = text_frame.add_paragraph()
                prompt_p.text = image_prompts[slide_index]
                prompt_p.alignment = PP_ALIGN.CENTER
                prompt_p.font.italic = True
                prompt_p.font.size = Pt(14)
                prompt_p.font.color.rgb = RGBColor(100, 100, 100)
                content_image_style = {
                    "left": image_position.get('left', 6.0),
                    "top": image_position.get('top', 2.0),
                    "width": image_position.get('width', 3.5),
                    "height": image_position.get('height', 4.0),
                    "fill_color": fill_color,
                    "border_color": border_color,
                    "border_width": image_slide_styles.get('border_width', 1.5),
                    "border_style": image_slide_styles.get('border_style', 'dashed')
                }
            
            preview_data["slides"].append({
                "type": "content",
                "title": slide_data.get("title", f"Slide {i+1}"),
                "points": slide_data.get("points", []),
                "points_styling": points_styling,
                "has_image": slide_index in image_prompts if image_prompts else False,
                "image_prompt": image_prompts.get(slide_index) if image_prompts else None,
                "image_style": content_image_style
            })
        
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_file.name)
        temp_file.close()
        return temp_file.name, preview_data
    
    except Exception as e:
        logger.error(f"PowerPoint creation error: {str(e)}")
        raise Exception(f"Failed to create PowerPoint: {str(e)}")
    

def generate_slide_previews(pptx_path):
    prs = Presentation(pptx_path)
    slide_previews = []
    
    for slide_idx, slide in enumerate(prs.slides):
        # Create a blank image for the slide (e.g., 960x540 pixels)
        img = Image.new('RGB', (960, 540), color=(255, 255, 255))
        # Here, you'd typically use a library like pptx2img or a custom renderer
        # For simplicity, we'll simulate with text rendering (upgrade with pptx2img later)
        from pptx.util import Inches
        from pptx.dml.color import RGBColor

        # Get slide layout (simplified)
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx + 1}"
        content = "\n".join([shape.text for shape in slide.shapes if shape.text and shape != slide.shapes.title])

        # Draw title
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.load_default()
            draw.text((50, 50), title, fill=(0, 0, 0), font=font)
            y_text = 100
            for line in content.split('\n'):
                draw.text((50, y_text), line, fill=(50, 50, 50), font=font)
                y_text += 30
        except Exception as e:
            print(f"Error rendering text for slide {slide_idx}: {e}")

        # Save image to a byte stream
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        slide_previews.append({
            'image': img_byte_arr,
            'title': title,
            'content': content
        })

    return slide_previews
    
@app.route('/generate_ppt', methods=['POST'])
@login_required
def generate_ppt():
    try:
        data = request.json
        template = data.get('template', 'default')
        content_type = data.get('content_type', 'auto_generate')  # 'auto_generate', 'custom'
        
        # Validate template
        if not template_manager.get_template(template):
            return jsonify({"error": "Invalid template selected"}), 400
        
        content_data = None
        topic = None
        image_prompts = {}
        
        # Handle different content types
        if content_type == 'auto_generate':
            # Original flow - generate from topic
            topic = data.get('topic')
            num_slides = int(data.get('num_slides', 3))
            
            if not topic:
                return jsonify({"error": "Topic is required for auto-generated content"}), 400
            if num_slides < 1 or num_slides > 20:
                return jsonify({"error": "Number of slides must be between 1 and 20"}), 400
                
            logger.info(f"User {session['username']} generating content for topic: {topic} with {num_slides} slides using template: {template}")
            content_data = generate_text_content(topic, num_slides)
            
        elif content_type == 'custom':
            # New flow - process custom content through Ollama
            custom_content = data.get('custom_content')
            custom_title = data.get('custom_title', 'Custom Presentation')
            
            if not custom_content:
                return jsonify({"error": "Custom content is required when selecting custom content type"}), 400
                
            logger.info(f"User {session['username']} using custom content with template: {template}")
            # Pass custom content to Ollama for processing
            content_data = generate_text_content(custom_title, 0, custom_content)
            topic = content_data.get("title", custom_title)
            
        else:
            return jsonify({"error": "Invalid content type"}), 400
        
        # Calculate slide count (title slide + content slides)
        slide_count = 1 + len(content_data.get('slides', []))
        
        # Generate image prompts for all slides
        try:
            logger.info("Generating title image prompt")
            title_image_prompt = generate_image_prompt(topic)
            if title_image_prompt:
                image_prompts["title"] = title_image_prompt
                
            for i, slide_data in enumerate(content_data.get("slides", [])):
                slide_title = slide_data.get("title", "")
                slide_image_prompt = generate_image_prompt(f"{topic} - {slide_title}")
                if slide_image_prompt:
                    image_prompts[str(i)] = slide_image_prompt
        except Exception as e:
            logger.warning(f"Image prompt generation failed: {str(e)}")
            
        # Create the presentation
        logger.info(f"Creating PowerPoint presentation with template: {template}")
        ppt_file, preview_data = create_presentation(content_data, image_prompts, template)
        
        # Save the file
        unique_id = uuid.uuid4().hex[:8]
        safe_topic = topic.replace(' ', '_') if topic else 'Presentation'
        filename = f"{safe_topic}_{unique_id}.pptx"
        user_filename = os.path.join("static", "downloads", filename)
        os.makedirs(os.path.dirname(user_filename), exist_ok=True)
        with open(ppt_file, 'rb') as src, open(user_filename, 'wb') as dst:
            dst.write(src.read())
        os.unlink(ppt_file)
        
        # Save presentation to user's history
        try:
            conn = get_db()
            c = conn.cursor()
            c.execute('''
                INSERT INTO presentations (user_id, title, filename, template, slide_count)
                VALUES (?, ?, ?, ?, ?)
            ''', (session['user_id'], content_data.get("title", topic), filename, template, slide_count))
            conn.commit()
            conn.close()
            logger.info(f"Saved presentation to user {session['username']}'s history")
        except sqlite3.Error as e:
            logger.error(f"Failed to save presentation to history: {str(e)}")
        
        return jsonify({
            "success": True,
            "filename": filename,
            "download_url": f"/static/downloads/{filename}",
            "content": content_data,
            "image_prompts": image_prompts,
            "template": template,
            "preview_data": preview_data
        })
    except Exception as e:
        logger.error(f"Error processing request: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/update_ppt', methods=['POST'])
@login_required
def update_ppt():
    try:
        data = request.json
        content_data = data.get('content')
        image_prompts = data.get('image_prompts', {})
        template = data.get('template', 'default')
        if not content_data or 'title' not in content_data or 'slides' not in content_data:
            return jsonify({"error": "Invalid presentation content"}), 400
        logger.info(f"User {session['username']} updating PowerPoint presentation")
        
        # Calculate slide count (title slide + content slides)
        slide_count = 1 + len(content_data.get('slides', []))
        
        ppt_file, preview_data = create_presentation(content_data, image_prompts, template)
        unique_id = uuid.uuid4().hex[:8]
        topic = content_data.get("title", "Presentation").replace(' ', '_')
        filename = f"{topic}_{unique_id}.pptx"
        user_filename = os.path.join("static", "downloads", filename)
        os.makedirs(os.path.dirname(user_filename), exist_ok=True)
        with open(ppt_file, 'rb') as src, open(user_filename, 'wb') as dst:
            dst.write(src.read())
        os.unlink(ppt_file)
        
        # Save updated presentation to user's history
        try:
            conn = get_db()
            c = conn.cursor()
            c.execute('''
                INSERT INTO presentations (user_id, title, filename, template, slide_count)
                VALUES (?, ?, ?, ?, ?)
            ''', (session['user_id'], content_data.get("title", topic), filename, template, slide_count))
            conn.commit()
            conn.close()
            logger.info(f"Saved updated presentation to user {session['username']}'s history")
        except sqlite3.Error as e:
            logger.error(f"Failed to save updated presentation to history: {str(e)}")
            
        return jsonify({
            "success": True,
            "filename": filename,
            "download_url": f"/static/downloads/{filename}",
            "preview_data": preview_data
        })
    except Exception as e:
        logger.error(f"Error updating presentation: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/welcome')
def welcome():
    if 'user_id' in session:
        return redirect(url_for('dashboard'))
    return render_template('welcome.html')

@app.route('/')
def index():
    if 'user_id' in session:
        return redirect(url_for('dashboard'))
    return redirect(url_for('welcome'))

from flask import flash, redirect, url_for, render_template, session
from datetime import datetime
from dateutil.relativedelta import relativedelta
import sqlite3

@app.route('/profile')
@login_required
def profile():
    print(f"Session contents: {session}")
    try:
        user_id = int(session['user_id'])
        print(f"User ID: {user_id}")
        conn = get_db()
        c = conn.cursor()
        
        # Fetch user data
        c.execute('SELECT username, email, created_at FROM users WHERE id = ?', (user_id,))
        user = c.fetchone()
        print(f"User: {user}")
        if not user:
            flash('User not found.', 'error')
            return redirect(url_for('login'))
        if not user['username']:
            flash('Invalid username.', 'error')
            return redirect(url_for('dashboard'))
        user = {'username': user['username'], 'email': user['email'], 'created_at': user['created_at']}
        
        # Fetch total presentation count
        c.execute('SELECT COUNT(*) as count FROM presentations WHERE user_id = ?', (user_id,))
        presentation_count = c.fetchone()['count']
        print(f"Presentation count: {presentation_count}")
        
        # Fetch monthly count
        month_start = datetime.now() - relativedelta(months=1)
        c.execute('SELECT COUNT(*) as count FROM presentations WHERE user_id = ? AND created_at >= ?', 
                  (user_id, month_start))
        presentations_this_month = c.fetchone()['count']
        print(f"Monthly count: {presentations_this_month}")
        
        # Fetch weekly count
        week_start = datetime.now() - relativedelta(days=7)
        c.execute('SELECT COUNT(*) as count FROM presentations WHERE user_id = ? AND created_at >= ?', 
                  (user_id, week_start))
        presentations_this_week = c.fetchone()['count']
        print(f"Weekly count: {presentations_this_week}")
        
        # Fetch graph data
        c.execute('SELECT DATE(created_at) as date, COUNT(*) as count FROM presentations WHERE user_id = ? AND created_at >= ? GROUP BY DATE(created_at)', 
                  (user_id, month_start))
        graph_data_raw = c.fetchall()
        graph_data = [{'date': row['date'], 'count': row['count']} for row in graph_data_raw]
        print(f"Graph data: {graph_data}")
        
        # Fetch recent presentations
        c.execute('SELECT title, created_at, slide_count FROM presentations WHERE user_id = ? ORDER BY created_at DESC LIMIT 5', 
                  (session['user_id'],))
        recent_presentations = c.fetchall()
        
        return render_template('profile.html', 
                              user=user, 
                              presentation_count=presentation_count,
                              presentations_this_month=presentations_this_month,
                              presentations_this_week=presentations_this_week,
                              graph_data=graph_data,
                              recent_presentations=recent_presentations)
    except Exception as e:
        print(f"Profile route error: {str(e)}")
        flash(f'Error: {str(e)}', 'error')
        return redirect(url_for('dashboard'))
    finally:
        conn.close()

@app.route('/static/<path:path>')
def serve_static(path):
    logger.debug(f"Serving static file: {path}")
    return app.send_static_file(path)

@app.route('/get_templates', methods=['GET'])
def get_templates():
    try:
        templates = template_manager.get_all_templates()
        template_response = {}
        for key, template in templates.items():
            template_response[key] = {
                "name": template.get('name', key),
                "description": template.get('description', ''),
                "preview_image": template.get('preview_image', ''),
                "styles": template.get('styles', {})
            }
        logger.info(f"Returning {len(template_response)} templates")
        return jsonify({
            "success": True,
            "templates": template_response
        })
    except Exception as e:
        logger.error(f"Error retrieving templates: {str(e)}")
        return jsonify({"error": "Failed to retrieve templates"}), 500

@app.route('/download/<filename>')
@login_required
def download_file(filename):
    file_path = os.path.join("static", "downloads", filename)
    if not os.path.exists(file_path):
        logger.error(f"Download file not found: {file_path}")
        return jsonify({"error": "File not found"}), 404
    
    # Verify this presentation belongs to the current user
    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT * FROM presentations WHERE user_id = ? AND filename = ?', 
             (session['user_id'], filename))
    presentation = c.fetchone()
    conn.close()
    
    if not presentation:
        logger.warning(f"User {session['username']} attempted to access unauthorized file: {filename}")
        return jsonify({"error": "Unauthorized access"}), 403
        
    logger.info(f"User {session['username']} downloading file: {file_path}")
    return send_file(file_path, as_attachment=True)

@app.route('/user/history')
@login_required
def user_history():
    conn = get_db()
    c = conn.cursor()
    c.execute('SELECT * FROM presentations WHERE user_id = ? ORDER BY created_at DESC', (session['user_id'],))
    presentations = c.fetchall()
    conn.close()
    
    return jsonify({
        "success": True,
        "history": [{
            "id": p['id'],
            "title": p['title'],
            "filename": p['filename'],
            "template": p['template'],
            "slide_count": p['slide_count'],
            "created_at": p['created_at'],
            "download_url": f"/download/{p['filename']}"
        } for p in presentations]
    })

if __name__ == '__main__':
    # Initialize database
    init_db()
    
    # Create directories if needed
    os.makedirs(os.path.join("static", "downloads"), exist_ok=True)
    
    app.run(debug=True)